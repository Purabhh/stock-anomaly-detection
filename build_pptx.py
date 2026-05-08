"""
Build presentation.pptx — the demo-video slideshow.

Run after editing if you want to regenerate:
    python build_pptx.py

Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# Colors
DARK_BLUE = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0xE7, 0x4C, 0x3C)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
CODE_BG = RGBColor(0x2D, 0x2D, 0x2D)
CODE_FG = RGBColor(0xE8, 0xE8, 0xE8)


def add_title_slide(prs, title, subtitle, footer):
    """First slide — big centered title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(60)
    r.font.bold = True
    r.font.color.rgb = DARK_BLUE
    # Subtitle
    sb = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(1.0))
    sf = sb.text_frame
    sp = sf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sr = sp.add_run()
    sr.text = subtitle
    sr.font.size = Pt(24)
    sr.font.color.rgb = DARK_GREY
    # Footer
    fb = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5))
    ff = fb.text_frame
    fp = ff.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER
    fr = fp.add_run()
    fr.text = footer
    fr.font.size = Pt(18)
    fr.font.color.rgb = DARK_GREY
    fr.font.italic = True
    return slide


def add_content_slide(prs, title, bullets=None, code=None, code_lang='python'):
    """Standard content slide: title + optional bullets + optional code."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title bar
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(12.5), Inches(0.7))
    tf = title_box.text_frame
    tf.margin_left = Inches(0.1)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = DARK_BLUE

    # Underline (a thin rectangle)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # Body — split between bullets and code if both present
    body_top = Inches(1.4)
    if bullets and code:
        # bullets on top half, code on bottom half
        bullet_box = slide.shapes.add_textbox(Inches(0.5), body_top, Inches(12.3), Inches(2.5))
        _fill_bullets(bullet_box, bullets)
        code_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.8))
        _fill_code(code_box, code)
    elif bullets:
        bullet_box = slide.shapes.add_textbox(Inches(0.5), body_top, Inches(12.3), Inches(5.5))
        _fill_bullets(bullet_box, bullets)
    elif code:
        code_box = slide.shapes.add_textbox(Inches(0.5), body_top, Inches(12.3), Inches(5.5))
        _fill_code(code_box, code)

    return slide


def _fill_bullets(box, bullets):
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        # Bullet char + text
        bullet_run = p.add_run()
        bullet_run.text = "•  "
        bullet_run.font.size = Pt(22)
        bullet_run.font.color.rgb = ACCENT
        bullet_run.font.bold = True

        text_run = p.add_run()
        text_run.text = b
        text_run.font.size = Pt(22)
        text_run.font.color.rgb = DARK_GREY


def _fill_code(box, code):
    """Render a code block in monospace on a dark background."""
    # Dark background rectangle
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    tf.word_wrap = True
    lines = code.strip('\n').split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        r = p.add_run()
        r.text = line if line else ' '
        r.font.name = 'Consolas'
        r.font.size = Pt(15)
        r.font.color.rgb = CODE_FG


def add_table_slide(prs, title, header, rows):
    """Slide that's a title + a data table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(12.5), Inches(0.7))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = DARK_BLUE
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # Table
    n_rows = len(rows) + 1
    n_cols = len(header)
    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.5), Inches(1.5),
        Inches(12.3), Inches(0.5 * n_rows + 0.5)
    )
    tbl = tbl_shape.table

    # Header
    for j, h in enumerate(header):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(16)
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.CENTER

    # Body rows
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(14)
                    r.font.color.rgb = DARK_GREY
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GREY

    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---------- Slide 1: Title ----------
    add_title_slide(
        prs,
        title="Anomalyy",
        subtitle="Multi-Method Detection and Semantic Classification\nof Stock-Market Anomalies",
        footer="CS 210: Data Management for Data Science  •  Purabh Singh"
    )

    # ---------- Slide 2: Problem ----------
    add_content_slide(
        prs, "The Problem",
        bullets=[
            "Markets generate massive daily data: prices, volumes, news, central-bank announcements.",
            "A small number of anomalous trading days drive an outsized share of risk and return.",
            "Detection alone is not enough; without context, a list of dates is not actionable.",
            "Research question: can unsupervised ML identify anomalies AND can external context explain them?",
        ]
    )

    # ---------- Slide 3: Pipeline Architecture ----------
    add_content_slide(
        prs, "Pipeline (6 Stages)",
        bullets=[
            "1.  Ingest: yfinance prices + GDELT BigQuery news sentiment.",
            "2.  Clean: schema validation, tz strip, V2Tone normalization.",
            "3.  Feature engineer: 40+ technical indicators per ticker.",
            "4.  Detect: Z-Score + Isolation Forest + LOF, agreement scoring.",
            "5.  Classify: 4-way precedence-based labeling.",
            "6.  Persist + visualize: SQLite database + Plotly per-ticker charts.",
        ]
    )

    # ---------- Slide 4: Data Sources ----------
    add_content_slide(
        prs, "Data Sources",
        bullets=[
            "yfinance: 5 tickers (AAPL, MSFT, TSLA, AMZN, ^GSPC), 2015 to 2024, 12,570 daily rows.",
            "GDELT 2.0 GKG via Google BigQuery: V2Tone sentiment + article counts (daily aggregates).",
            "FOMC dates: 78 hardcoded Fed meetings (2015 to 2024).",
            "Why GDELT BQ? NewsAPI free tier caps at 30 days; GDELT BQ covers full 10-year history.",
        ],
        code='''# news_bq.py — daily aggregate query
SELECT
  DATE(_PARTITIONTIME) AS date,
  AVG(SAFE_CAST(SPLIT(V2Tone, ',')[OFFSET(0)] AS FLOAT64)) / 100.0 AS sentiment_compound,
  COUNT(*) AS article_count
FROM `gdelt-bq.gdeltv2.gkg_partitioned`
WHERE UPPER(V2Organizations) LIKE '%APPLE INC%'
GROUP BY date'''
    )

    # ---------- Slide 5: Database Schema ----------
    add_content_slide(
        prs, "Database Schema (4 Normalized Tables)",
        bullets=[
            "stocks (PK: symbol) — company metadata, stored once per ticker.",
            "price_data (FK: symbol, UNIQUE: symbol+date) — daily OHLCV rows.",
            "news_articles (FK: symbol, UNIQUE: url) — daily V2Tone aggregates + article_count.",
            "anomalies (FK: symbol, UNIQUE: symbol+anomaly_date) — flags, agreement, label.",
            "Indexes on (symbol, date) speed up the per-ticker date-range queries.",
        ]
    )

    # ---------- Slide 6: Feature Engineering ----------
    add_content_slide(
        prs, "Feature Engineering — 40+ Indicators",
        bullets=[
            "Returns: daily, log, ratio-based.",
            "Moving averages: SMA + EMA at 5, 10, 20, 50, 200-day windows.",
            "Volatility: ATR, rolling std, Bollinger Bands, drawdown.",
            "Momentum: RSI-14/28, MACD, MACD histogram.",
            "Volume + statistical (OBV, VWAP, skew, kurtosis, Hurst) + candle patterns.",
        ],
        code='''from src.feature_engineering import FeatureEngineer
engineer = FeatureEngineer()
features = engineer.engineer_features(price_df)
# features.shape -> (~2300, 50+)'''
    )

    # ---------- Slide 7: Three Detection Methods ----------
    add_content_slide(
        prs, "Three Complementary Detectors",
        bullets=[
            "Z-Score (univariate, 3σ): catches extreme single-feature spikes.",
            "Isolation Forest (multivariate, contamination=10%): isolates rare points by random splits.",
            "Local Outlier Factor (density-based, n=20): detects locally rare points.",
            "Each method targets a different kind of 'unusual' — consensus is more robust.",
        ],
        code='''detector = AnomalyDetector(contamination=0.1)
results = detector.detect_anomalies(features)
# columns added: z_score_anomaly, isolation_forest_anomaly,
#                lof_anomaly, agreement_score (0-3)'''
    )

    # ---------- Slide 8: Agreement + Classification ----------
    add_content_slide(
        prs, "Agreement Scoring + 4-Way Classification",
        bullets=[
            "agreement_score = sum of method flags (0 to 3).",
            "Threshold ≥ 2 keeps consensus only: 1,573 days = ~12% of all trading days.",
            "Precedence: macroeconomic_event → sector_contagion → vader_sentiment_spike → unexplained.",
            "Sentiment rule uses per-ticker z-score, NOT absolute threshold (V2Tone clusters near 0).",
        ],
        code='''# classify_anomaly_type — strict precedence
if abs(date - fomc_date).days <= 2: return 'macroeconomic_event'
if date in contagion_dates: return 'sector_contagion'
if news_volume >= 5 and abs(z) >= 2.0: return 'vader_sentiment_spike'
return 'unexplained' '''
    )

    # ---------- Slide 9: Pivots from Proposal ----------
    add_content_slide(
        prs, "Pivots From the February Proposal",
        bullets=[
            "News source: NewsAPI (30-day cap) → GDELT DOC API (throttled) → GDELT BigQuery (full history).",
            "Sentiment: VADER on headlines → GDELT V2Tone on full article bodies (VADER kept as fallback).",
            "Classification: manual labeling → automatic precedence rules (1,573 anomalies = manual is impractical).",
            "Categories: 'earnings' replaced with 'vader_sentiment_spike' (no earnings calendar in scope).",
        ]
    )

    # ---------- Slide 10: Results ----------
    add_table_slide(
        prs, "Results — 1,573 High-Confidence Anomalies",
        header=["Ticker", "Total", "macro", "contagion", "sentiment", "unexplained"],
        rows=[
            ["AAPL", 325, 72, 170, 3, 80],
            ["MSFT", 338, 56, 183, 1, 98],
            ["TSLA", 309, 56, 167, 3, 83],
            ["AMZN", 293, 47, 146, 0, 100],
            ["^GSPC", 308, 53, 175, 2, 78],
            ["Total", 1573, 284, 841, 9, 439],
        ]
    )

    # ---------- Slide 11: Validation ----------
    add_content_slide(
        prs, "Validation: 2× Lift Over Random Baseline",
        bullets=[
            "Precision-by-explanation = (284 + 841 + 9) / 1,573 = 72.1%.",
            "Roughly 7 of 10 high-confidence anomalies map to a known event.",
            "Random FOMC overlap baseline: 78 × 3 / 2,514 ≈ 9.3%.",
            "Observed FOMC overlap: 18.1% — roughly 2× random — real signal, not chance.",
            "pytest: 8/8 passing across feature engineering, detection, schema, classification.",
        ]
    )

    # ---------- Slide 12: Limitations + Future Work ----------
    add_content_slide(
        prs, "Limitations & Future Work",
        bullets=[
            "No earnings calendar — likely accounts for most of the 28% unexplained bucket.",
            "V2Tone is corpus-aggregated — single outlier articles get smoothed away.",
            "contamination=0.1 fixed — directly determines anomaly-set size, no formal tuning.",
            "5 tickers is small + tech-tilted — limits sector-level pattern testing.",
            "Future: earnings calendar, per-article V2Tone, walk-forward baselines, GARCH, 50+ tickers.",
        ]
    )

    # ---------- Slide 13: Conclusion ----------
    add_content_slide(
        prs, "Conclusion",
        bullets=[
            "Detection works: three-method consensus identifies ~12% of days, well above noise floor.",
            "Explanation works: 72.1% of anomalies map to a specific known cause.",
            "Combined output is actionable: 'Mar 16, 2020: sector_contagion' >> 'Mar 16, 2020: anomaly'.",
            "Source: github.com/Purabhh/Anomalyy   •   Thank you.",
        ]
    )

    out = 'presentation.pptx'
    prs.save(out)
    import os
    print(f'Wrote {out} ({os.path.getsize(out)} bytes, {len(prs.slides)} slides)')


if __name__ == '__main__':
    build()
